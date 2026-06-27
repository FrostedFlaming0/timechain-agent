"""
extractors.py — text extraction for ingested content.

Resurrected from the removed file_ingest pipeline: ingest_blob routes
uploads through here so PDFs, Office documents, spreadsheets, slide decks,
and delimited files keep their TEXT searchable (sealed on the record,
embedded, retrievable) instead of silently becoming opaque blobs.

Every format-specific dependency is OPTIONAL: a missing library degrades to
a descriptive placeholder (extraction_method="missing-dep"), never an error.
Install any of: pypdf, python-docx, openpyxl, python-pptx, Pillow, chardet.
"""

from __future__ import annotations

import csv
import io
from pathlib import Path

# Cap on extracted text per file — large enough for real documents, small
# enough that one record can never balloon the chain. Callers store the
# returned `truncated` flag so the cut is visible, not silent.
MAX_EXTRACTED_CHARS = 40_000

# Extensions safe to decode as plain text when the MIME type is not text/*.
_TEXT_EXTS = {
    ".txt", ".md", ".rst", ".log", ".json", ".jsonl", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".xml", ".html", ".htm", ".py", ".js", ".ts",
    ".c", ".h", ".cpp", ".hpp", ".java", ".go", ".rs", ".rb", ".sh",
    ".sql", ".css",
}


def extract_text(raw: bytes, filename: str,
                 mime_type: str = "") -> tuple[str, str, bool]:
    """
    Extract searchable text from `raw`. Returns
    `(text, extraction_method, truncated)`.

    Dispatch is by filename extension first, MIME type second. Formats with
    no extractor (unknown binary) return an empty string with method
    "none" — the caller decides how to describe the blob.
    """
    ext = Path(filename).suffix.lower()
    try:
        if mime_type.startswith("image/"):
            text, method = _extract_image_metadata(raw, ext, filename)
        elif ext == ".pdf" or mime_type == "application/pdf":
            text, method = _extract_pdf(raw)
        elif ext in (".docx", ".dotx"):
            text, method = _extract_docx(raw)
        elif ext == ".xlsx":
            text, method = _extract_xlsx(raw)
        elif ext == ".csv":
            text, method = _extract_csv(raw, delimiter=",")
        elif ext == ".tsv":
            text, method = _extract_csv(raw, delimiter="\t")
        elif ext == ".pptx":
            text, method = _extract_pptx(raw)
        elif mime_type.startswith("text/") or ext in _TEXT_EXTS:
            text, method = _extract_plaintext(raw)
        else:
            # Unknown binary: decoding it as text would seal garbage.
            text, method = "", "none"
    except Exception as e:                       # noqa: BLE001 — extraction
        # must never sink an ingest; the record still carries metadata.
        text = f"[extraction failed for {filename}: {type(e).__name__}: {e}]"
        method = "failed"

    truncated = False
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n\n[... truncated ...]"
        truncated = True
    return text, method, truncated


def _extract_plaintext(raw: bytes) -> tuple[str, str]:
    # Try utf-8 first, then fall back to chardet detection if available.
    try:
        return raw.decode("utf-8"), "utf8"
    except UnicodeDecodeError:
        pass
    try:
        import chardet
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "latin-1"
        return raw.decode(encoding, errors="replace"), f"chardet-{encoding}"
    except ImportError:
        return raw.decode("latin-1", errors="replace"), "latin1-fallback"


def _extract_pdf(raw: bytes) -> tuple[str, str]:
    try:
        import pypdf
    except ImportError:
        return "[pdf extraction needs `pip install pypdf`]", "missing-dep"
    reader = pypdf.PdfReader(io.BytesIO(raw))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(f"--- page {i + 1} ---\n{page.extract_text() or ''}")
        except Exception as e:                   # noqa: BLE001
            pages.append(f"--- page {i + 1} ---\n[error extracting: {e}]")
    return "\n\n".join(pages).strip(), "pypdf"


def _extract_docx(raw: bytes) -> tuple[str, str]:
    try:
        import docx
    except ImportError:
        return ("[docx extraction needs `pip install python-docx`]",
                "missing-dep")
    doc = docx.Document(io.BytesIO(raw))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables.
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts), "python-docx"


def _extract_xlsx(raw: bytes) -> tuple[str, str]:
    try:
        import openpyxl
    except ImportError:
        return "[xlsx extraction needs `pip install openpyxl`]", "missing-dep"
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True,
                                read_only=True)
    out_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out_parts.append(f"--- sheet: {sheet_name} ---")
        rows_rendered = 0
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cell.strip() for cell in cells):
                out_parts.append(" | ".join(cells))
                rows_rendered += 1
            if rows_rendered >= 1000:
                out_parts.append("[... more rows truncated ...]")
                break
    wb.close()
    return "\n".join(out_parts), "openpyxl"


def _extract_csv(raw: bytes, delimiter: str = ",") -> tuple[str, str]:
    text, _ = _extract_plaintext(raw)
    # Re-render through the csv module to normalize quoting and produce a
    # tabular, line-oriented view that's easier for the LLM to read.
    out_lines = []
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    for i, row in enumerate(reader):
        out_lines.append(" | ".join(row))
        if i >= 2000:
            out_lines.append("[... more rows truncated ...]")
            break
    return "\n".join(out_lines), "csv"


def _extract_pptx(raw: bytes) -> tuple[str, str]:
    try:
        import pptx
    except ImportError:
        return ("[pptx extraction needs `pip install python-pptx`]",
                "missing-dep")
    prs = pptx.Presentation(io.BytesIO(raw))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"--- slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs)
                    if txt.strip():
                        slide_lines.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"[notes] {notes}")
        parts.append("\n".join(slide_lines))
    return "\n\n".join(parts), "python-pptx"


def _extract_image_metadata(raw: bytes, ext: str,
                            filename: str) -> tuple[str, str]:
    """
    Images are not OCR'd or captioned here — just described. The actual
    visual content is sent to vision-capable LLMs separately when the
    record is in retrieval context (see agent._collect_attachments).
    """
    try:
        from PIL import Image
    except ImportError:
        return (f"[image: {filename} ({len(raw):,} bytes, {ext})]",
                "missing-dep")
    try:
        # `with` releases the underlying buffer deterministically.
        with Image.open(io.BytesIO(raw)) as img:
            return (
                f"[image: {filename} ({img.width}x{img.height} {img.mode}, "
                f"format={img.format}, {len(raw):,} bytes)]"
            ), "pillow-metadata"
    except Exception as e:                       # noqa: BLE001
        return f"[image: {filename} (could not parse: {e})]", "pillow-failed"
