"""
file_ingest — read user files and produce normalized chain content.

Handles the file types the project supports:
  - Documents/text: .pdf, .doc, .docx, .dot, .dotx, .txt, .rtf, .md, .hwp, .hwpx
  - Data/spreadsheets: .xlsx, .xls, .csv, .tsv
  - Presentations: .pptx, .ppt
  - Images: .jpg, .jpeg, .png, .webp, .heic, .heif, .gif, .bmp
  - Code: .json, .yaml/.yml, .xml, .html, .css, .js, .ts, .py, .sh,
    .c, .cpp, .h, .java, .rs, .go, plus other plain-text source files

Two outputs per file:
  1. A blob stored in `blobs/<sha256>` — the raw bytes, content-addressed.
  2. A 'file' record on the chain containing metadata + extracted text:
        {
            "filename": "report.pdf",
            "ext": ".pdf",
            "kind": "document" | "image" | "code" | "spreadsheet" | "presentation",
            "size_bytes": 123456,
            "blob_sha256": "deadbeef...",
            "extracted_text": "...",       # text representation for LLM context
            "extraction_method": "pypdf",  # how the text was produced
            "extraction_truncated": False, # True if extracted_text was cut short
        }

The chain record is signed and tamper-evident the same as any other.
The blob is content-addressed by its sha256 — if the blob is altered,
the hash on the record won't match, and that's detectable by checking
sha256(blob) == record.content["blob_sha256"].

Optional dependencies (install only what you need):
  pip install pypdf python-docx openpyxl python-pptx Pillow chardet
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Optional


# Maximum extracted text length stored in a record (chars). Longer files
# get truncated with a marker. Tunable; 40K chars is roughly 10K tokens.
MAX_EXTRACTED_CHARS = 40_000

# Maximum bytes we'll ingest per file. Prevents accidental huge uploads
# from blowing up the chain or blob store. Override per call if needed.
DEFAULT_MAX_FILE_BYTES = 50 * 1024 * 1024  # 50 MB


# Extension classification — drives both extraction strategy and the
# 'kind' field on the record.
DOCUMENT_EXTS = {
    ".pdf", ".doc", ".docx", ".dot", ".dotx", ".txt", ".rtf",
    ".md", ".markdown", ".hwp", ".hwpx", ".odt",
}
SPREADSHEET_EXTS = {".xlsx", ".xls", ".csv", ".tsv", ".ods"}
PRESENTATION_EXTS = {".pptx", ".ppt", ".odp"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".heic", ".heif",
              ".gif", ".bmp", ".tiff", ".tif"}
CODE_EXTS = {
    ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".css",
    ".js", ".ts", ".jsx", ".tsx", ".py", ".sh", ".bash", ".zsh",
    ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".java", ".rs",
    ".go", ".rb", ".php", ".swift", ".kt", ".scala", ".lua",
    ".sql", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".dockerfile", ".makefile",
}

# Anything plain-text that doesn't fit the buckets above
TEXT_EXTS = {".log"}

ALL_SUPPORTED = (
    DOCUMENT_EXTS | SPREADSHEET_EXTS | PRESENTATION_EXTS
    | IMAGE_EXTS | CODE_EXTS | TEXT_EXTS
)


@dataclass
class IngestResult:
    """What ingest_file returns. Caller turns this into a chain record."""
    filename: str
    ext: str
    kind: str
    size_bytes: int
    blob_sha256: str
    blob_path: str  # relative to data dir
    extracted_text: str
    extraction_method: str
    extraction_truncated: bool

    def to_record_content(self) -> dict:
        """Shape suitable for chain.append('file', content=this)."""
        return {
            "filename": self.filename,
            "ext": self.ext,
            "kind": self.kind,
            "size_bytes": self.size_bytes,
            "blob_sha256": self.blob_sha256,
            "blob_path": self.blob_path,
            "extracted_text": self.extracted_text,
            "extraction_method": self.extraction_method,
            "extraction_truncated": self.extraction_truncated,
            "schema_version": 1,
        }


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def is_supported(path: str | Path) -> bool:
    return Path(path).suffix.lower() in ALL_SUPPORTED


def classify_kind(ext: str) -> str:
    ext = ext.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in SPREADSHEET_EXTS:
        return "spreadsheet"
    if ext in PRESENTATION_EXTS:
        return "presentation"
    if ext in CODE_EXTS:
        return "code"
    if ext in DOCUMENT_EXTS:
        return "document"
    return "text"


def ingest_file(
    path: str | Path,
    blob_dir: str | Path,
    max_bytes: int = DEFAULT_MAX_FILE_BYTES,
) -> IngestResult:
    """
    Read a file, store its bytes content-addressed in `blob_dir`, and
    return an IngestResult ready to be turned into a chain record.

    Raises FileNotFoundError, ValueError (unsupported), or OSError (too big).
    """
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        raise FileNotFoundError(p)

    ext = p.suffix.lower()
    if ext not in ALL_SUPPORTED:
        raise ValueError(f"unsupported file type: {ext}")

    size = p.stat().st_size
    if size > max_bytes:
        raise OSError(
            f"file too large: {size:,} bytes (limit {max_bytes:,}). "
            "Override with max_bytes if you really want to ingest this."
        )

    raw = p.read_bytes()
    sha = hashlib.sha256(raw).hexdigest()

    blob_dir = Path(blob_dir)
    blob_dir.mkdir(parents=True, exist_ok=True)
    blob_path = blob_dir / sha
    if not blob_path.exists():
        blob_path.write_bytes(raw)

    kind = classify_kind(ext)
    text, method, truncated = _extract_text(raw, ext, kind, p.name)

    return IngestResult(
        filename=p.name,
        ext=ext,
        kind=kind,
        size_bytes=size,
        blob_sha256=sha,
        blob_path=str(blob_path.name),  # store just the basename; resolve via blob_dir
        extracted_text=text,
        extraction_method=method,
        extraction_truncated=truncated,
    )


def verify_blob(content: dict, blob_dir: str | Path) -> bool:
    """
    Confirm the on-disk blob matches the sha256 recorded on the chain.
    Returns True if the blob is intact, False if missing or modified.
    """
    blob_path = Path(blob_dir) / content["blob_path"]
    if not blob_path.exists():
        return False
    actual = hashlib.sha256(blob_path.read_bytes()).hexdigest()
    return actual == content["blob_sha256"]


# ---------------------------------------------------------------------------
# Text extraction — one function per kind, tried in order
# ---------------------------------------------------------------------------

def _extract_text(raw: bytes, ext: str, kind: str, filename: str) -> tuple[str, str, bool]:
    """
    Return (text, extraction_method, truncated). If a specific extractor
    isn't available (missing dep, unsupported variant), falls back to a
    descriptive placeholder so the record still carries useful metadata.
    """
    try:
        if kind == "image":
            text, method = _extract_image_metadata(raw, ext, filename)
        elif ext == ".pdf":
            text, method = _extract_pdf(raw)
        elif ext in (".docx", ".dotx"):
            text, method = _extract_docx(raw)
        elif ext == ".xlsx":
            text, method = _extract_xlsx(raw)
        elif ext in (".csv",):
            text, method = _extract_csv(raw, delimiter=",")
        elif ext == ".tsv":
            text, method = _extract_csv(raw, delimiter="\t")
        elif ext == ".pptx":
            text, method = _extract_pptx(raw)
        elif ext in DOCUMENT_EXTS or ext in CODE_EXTS or ext in TEXT_EXTS:
            text, method = _extract_plaintext(raw)
        else:
            text, method = _extract_plaintext(raw)
    except Exception as e:
        text = f"[extraction failed for {filename}: {type(e).__name__}: {e}]"
        method = "failed"

    truncated = False
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n\n[... truncated ...]"
        truncated = True
    return text, method, truncated


def _extract_plaintext(raw: bytes) -> tuple[str, str]:
    # Try utf-8 first, then fall back to chardet detection if available
    try:
        return raw.decode("utf-8"), "utf8"
    except UnicodeDecodeError:
        pass
    try:
        import chardet
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "latin-1"
        return raw.decode(encoding, errors="replace"), f"chardet-{encoding}"
    except ImportError:
        return raw.decode("latin-1", errors="replace"), "latin1-fallback"


def _extract_pdf(raw: bytes) -> tuple[str, str]:
    try:
        import pypdf
    except ImportError:
        return "[pdf extraction needs `pip install pypdf`]", "missing-dep"
    reader = pypdf.PdfReader(io.BytesIO(raw))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(f"--- page {i + 1} ---\n{page.extract_text() or ''}")
        except Exception as e:
            pages.append(f"--- page {i + 1} ---\n[error extracting: {e}]")
    return "\n\n".join(pages).strip(), "pypdf"


def _extract_docx(raw: bytes) -> tuple[str, str]:
    try:
        import docx
    except ImportError:
        return "[docx extraction needs `pip install python-docx`]", "missing-dep"
    doc = docx.Document(io.BytesIO(raw))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts), "python-docx"


def _extract_xlsx(raw: bytes) -> tuple[str, str]:
    try:
        import openpyxl
    except ImportError:
        return "[xlsx extraction needs `pip install openpyxl`]", "missing-dep"
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    out_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out_parts.append(f"--- sheet: {sheet_name} ---")
        rows_rendered = 0
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cell.strip() for cell in cells):
                out_parts.append(" | ".join(cells))
                rows_rendered += 1
            if rows_rendered >= 1000:
                out_parts.append("[... more rows truncated ...]")
                break
    wb.close()
    return "\n".join(out_parts), "openpyxl"


def _extract_csv(raw: bytes, delimiter: str = ",") -> tuple[str, str]:
    text, _ = _extract_plaintext(raw)
    # Re-render through csv module to normalize quoting and produce
    # a tabular, line-oriented view that's easier for the LLM to read.
    out_lines = []
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    for i, row in enumerate(reader):
        out_lines.append(" | ".join(row))
        if i >= 2000:
            out_lines.append("[... more rows truncated ...]")
            break
    return "\n".join(out_lines), "csv"


def _extract_pptx(raw: bytes) -> tuple[str, str]:
    try:
        import pptx
    except ImportError:
        return "[pptx extraction needs `pip install python-pptx`]", "missing-dep"
    prs = pptx.Presentation(io.BytesIO(raw))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"--- slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs)
                    if txt.strip():
                        slide_lines.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"[notes] {notes}")
        parts.append("\n".join(slide_lines))
    return "\n\n".join(parts), "python-pptx"


def _extract_image_metadata(raw: bytes, ext: str, filename: str) -> tuple[str, str]:
    """
    For images we don't OCR or caption here — we just record metadata.
    The actual visual content gets sent to vision-capable LLMs separately
    when the image record is in retrieval context (see llm_clients.py).
    """
    try:
        from PIL import Image
    except ImportError:
        return f"[image: {filename} ({len(raw):,} bytes, {ext})]", "missing-dep"
    try:
        img = Image.open(io.BytesIO(raw))
        return (
            f"[image: {filename} ({img.width}x{img.height} {img.mode}, "
            f"format={img.format}, {len(raw):,} bytes)]"
        ), "pillow-metadata"
    except Exception as e:
        return f"[image: {filename} (could not parse: {e})]", "pillow-failed"
